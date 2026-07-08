"""Рендер презентации в PPTX по фирменному шаблону (план §5.1, M4).
Двуязычные слайды: uz и ru рендерятся в одном файле (заголовок uz / ru)."""
import io

from pptx import Presentation as PptxPresentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Pt

from app.core import storage

_SLIDE_W = Emu(9144000)   # 16:9 ширина
_SLIDE_H = Emu(5143500)   # 16:9 высота


def _hex(color: str) -> RGBColor:
    return RGBColor.from_string(color.lstrip("#").upper())


def _lang_title(slide: dict, lang: str) -> str:
    return slide.get(f"title_{lang}", "") or ""


def _lang_bullets(slide: dict, lang: str) -> list[str]:
    return slide.get(f"bullets_{lang}", []) or []


def _lang_notes(slide: dict, lang: str) -> str:
    return slide.get(f"notes_{lang}", "") or ""


def render_pptx(slides: list[dict], template: dict | None, lang: str = "ru") -> bytes:
    primary = _hex((template or {}).get("primary_color", "0D9488"))
    accent = _hex((template or {}).get("accent_color", "0F172A"))
    logo_url = (template or {}).get("logo_url")

    prs = PptxPresentation()
    prs.slide_width = _SLIDE_W
    prs.slide_height = _SLIDE_H
    blank = prs.slide_layouts[6]

    logo_bytes = None
    if logo_url:
        try:
            logo_bytes = storage.full_path(logo_url).read_bytes()
        except OSError:
            logo_bytes = None

    for index, slide_data in enumerate(slides):
        slide = prs.slides.add_slide(blank)
        title = _lang_title(slide_data, lang)
        bullets = _lang_bullets(slide_data, lang)
        is_title_slide = index == 0

        # Фон-акцент сверху
        bar = slide.shapes.add_shape(1, 0, 0, _SLIDE_W, Emu(400000))
        bar.fill.solid()
        bar.fill.fore_color.rgb = primary
        bar.line.fill.background()
        bar.shadow.inherit = False

        if is_title_slide:
            box = slide.shapes.add_textbox(Emu(700000), Emu(1800000), Emu(7700000), Emu(1500000))
            tf = box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = title
            run.font.size = Pt(40)
            run.font.bold = True
            run.font.color.rgb = accent
        else:
            # Заголовок
            title_box = slide.shapes.add_textbox(Emu(600000), Emu(550000), Emu(7900000), Emu(900000))
            tf = title_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = title
            run.font.size = Pt(28)
            run.font.bold = True
            run.font.color.rgb = primary

            # Есть картинка? делим слайд: текст слева, картинка справа
            image_url = slide_data.get("image_url")
            image_bytes = None
            if image_url:
                try:
                    image_bytes = storage.full_path(image_url).read_bytes()
                except OSError:
                    image_bytes = None

            text_width = Emu(4600000) if image_bytes else Emu(7900000)
            body_box = slide.shapes.add_textbox(Emu(600000), Emu(1600000), text_width, Emu(3100000))
            body = body_box.text_frame
            body.word_wrap = True
            for i, bullet in enumerate(bullets):
                p = body.paragraphs[0] if i == 0 else body.add_paragraph()
                run = p.add_run()
                run.text = f"•  {bullet}"
                run.font.size = Pt(18)
                run.font.color.rgb = accent
                p.space_after = Pt(10)

            if image_bytes:
                try:
                    slide.shapes.add_picture(
                        io.BytesIO(image_bytes), Emu(5400000), Emu(1600000),
                        width=Emu(3200000), height=Emu(2400000),
                    )
                except Exception:
                    pass  # битую картинку просто пропускаем

        # Логотип в углу
        if logo_bytes:
            try:
                slide.shapes.add_picture(io.BytesIO(logo_bytes), Emu(8100000), Emu(60000), height=Emu(300000))
            except Exception:
                pass

        # Заметки докладчика
        notes = _lang_notes(slide_data, lang)
        if notes:
            slide.notes_slide.notes_text_frame.text = notes

    buffer = io.BytesIO()
    prs.save(buffer)
    return buffer.getvalue()
