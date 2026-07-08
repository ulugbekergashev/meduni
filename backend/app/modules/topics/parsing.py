"""Извлечение текста из загруженных материалов (dev-вариант parse_worker).
OCR для сканов — later (vision-LLM), см. план §5.1."""
import io

SUPPORTED = {"pdf", "docx", "pptx", "txt", "md"}


def extract_text(data: bytes, file_type: str) -> str:
    if file_type in ("txt", "md"):
        return data.decode("utf-8", errors="ignore")

    if file_type == "pdf":
        from pypdf import PdfReader

        reader = PdfReader(io.BytesIO(data))
        return "\n".join(page.extract_text() or "" for page in reader.pages)

    if file_type == "docx":
        import docx

        document = docx.Document(io.BytesIO(data))
        parts = [p.text for p in document.paragraphs]
        for table in document.tables:
            for row in table.rows:
                parts.append(" | ".join(cell.text for cell in row.cells))
        return "\n".join(parts)

    if file_type == "pptx":
        from pptx import Presentation

        presentation = Presentation(io.BytesIO(data))
        parts = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    parts.append(shape.text_frame.text)
        return "\n".join(parts)

    raise ValueError(f"Неподдерживаемый формат: {file_type}")
